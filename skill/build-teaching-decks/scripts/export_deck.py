#!/usr/bin/env python3
"""Export ordered SVG slides through an editable backend or an explicit raster fallback."""

from __future__ import annotations

import argparse
import importlib.util
import json
import re
import shutil
import subprocess
import sys
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from xml.etree import ElementTree as ET

from runtime_detection import EXECUTORS, detect_executor, find_ppt_master

SEQ_RE = re.compile(r"^[A-Za-z_-]*?(\d+)(?:[_-]|$)")
NUMBER_RE = re.compile(r"^-?\d+(?:\.\d+)?")


def ordered_svgs(project: Path) -> list[Path]:
    files = list((project / "svg_output").glob("*.svg"))
    parsed: list[tuple[int, str, Path]] = []
    for path in files:
        match = SEQ_RE.match(path.stem)
        if not match:
            raise ValueError(f"SVG lacks a numeric sequence prefix: {path.name}")
        parsed.append((int(match.group(1)), path.name, path))
    parsed.sort()
    numbers = [item[0] for item in parsed]
    if numbers and numbers != list(range(numbers[0], numbers[-1] + 1)):
        raise ValueError("SVG sequence is not contiguous")
    lexical = sorted(files, key=lambda item: item.name)
    numeric = [item[2] for item in parsed]
    if lexical != numeric:
        raise ValueError("Lexicographic SVG order differs from numeric order; zero-pad filenames")
    return numeric


def svg_size(path: Path) -> tuple[float, float]:
    root = ET.parse(path).getroot()
    viewbox = root.get("viewBox")
    if viewbox:
        parts = re.split(r"[\s,]+", viewbox.strip())
        if len(parts) == 4:
            return float(parts[2]), float(parts[3])
    def parse(value: str | None) -> float | None:
        if not value:
            return None
        match = NUMBER_RE.match(value)
        return float(match.group()) if match else None
    width, height = parse(root.get("width")), parse(root.get("height"))
    if width and height:
        return width, height
    raise ValueError(f"Cannot determine SVG canvas size: {path}")


def detect_renderer() -> str | None:
    if importlib.util.find_spec("cairosvg"):
        return "cairosvg"
    if shutil.which("inkscape"):
        return "inkscape"
    if shutil.which("rsvg-convert"):
        return "rsvg-convert"
    if importlib.util.find_spec("svglib") and importlib.util.find_spec("reportlab"):
        return "svglib-reportlab"
    return None


def render_svg(source: Path, target: Path, renderer: str, width: int, height: int) -> None:
    target.parent.mkdir(parents=True, exist_ok=True)
    if renderer == "cairosvg":
        import cairosvg

        cairosvg.svg2png(
            url=str(source),
            write_to=str(target),
            output_width=width,
            output_height=height,
        )
        return
    if renderer == "inkscape":
        subprocess.run(
            [
                shutil.which("inkscape") or "inkscape",
                str(source),
                "--export-type=png",
                f"--export-filename={target}",
                f"--export-width={width}",
                f"--export-height={height}",
            ],
            check=True,
        )
        return
    if renderer == "rsvg-convert":
        subprocess.run(
            [
                shutil.which("rsvg-convert") or "rsvg-convert",
                "-w",
                str(width),
                "-h",
                str(height),
                "-o",
                str(target),
                str(source),
            ],
            check=True,
        )
        return
    if renderer == "svglib-reportlab":
        from reportlab.graphics import renderPM
        from svglib.svglib import svg2rlg

        drawing = svg2rlg(str(source))
        if drawing is None:
            raise RuntimeError(f"svglib could not parse {source}")
        scale_x = width / drawing.width
        scale_y = height / drawing.height
        drawing.scale(scale_x, scale_y)
        drawing.width = width
        drawing.height = height
        renderPM.drawToFile(drawing, str(target), fmt="PNG")
        return
    raise RuntimeError(f"Unknown renderer: {renderer}")


def export_raster(project: Path, output: Path, svgs: list[Path]) -> dict[str, object]:
    if importlib.util.find_spec("pptx") is None:
        raise RuntimeError("python-pptx is required for raster fallback export")
    renderer = detect_renderer()
    if not renderer:
        raise RuntimeError("No SVG renderer found; install cairosvg or provide Inkscape/rsvg-convert")
    from pptx import Presentation
    from pptx.util import Inches

    width_px, height_px = svg_size(svgs[0])
    prs = Presentation()
    prs.slide_width = Inches(width_px / 96)
    prs.slide_height = Inches(height_px / 96)
    blank = prs.slide_layouts[6]
    with tempfile.TemporaryDirectory(prefix="teaching-deck-render-") as temp:
        temp_dir = Path(temp)
        for index, svg in enumerate(svgs, start=1):
            png = temp_dir / f"{index:04d}.png"
            render_svg(svg, png, renderer, int(width_px * 2), int(height_px * 2))
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                str(png),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
        prs.core_properties.title = project.name
        prs.core_properties.comments = "Raster fallback export; slide elements are not editable."
        prs.save(output)
    return {
        "backend": f"raster-{renderer}",
        "capability_class": "raster-fallback",
        "editable": False,
        "slide_count": len(svgs),
        "canvas_px": [width_px, height_px],
    }


def export_ppt_master(project: Path, output: Path, root: Path) -> dict[str, object]:
    command = [
        sys.executable,
        str(root / "scripts" / "svg_to_pptx.py"),
        str(project),
        "-o",
        str(output),
        "--only",
        "native",
        "-s",
        "output",
    ]
    subprocess.run(command, check=True)
    return {
        "backend": "ppt-master-native",
        "capability_class": "editable-native",
        "editable": True,
        "command": command,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("project", type=Path, help="Deck project containing svg_output/")
    parser.add_argument("-o", "--output", type=Path, required=True)
    parser.add_argument("--backend", choices=("auto", "ppt-master", "raster"), default="auto")
    parser.add_argument(
        "--executor",
        choices=EXECUTORS,
        default="auto",
        help="Override executor detection; TEACHING_DECK_EXECUTOR is also supported",
    )
    parser.add_argument("--require-editable", action="store_true")
    parser.add_argument("--report", type=Path)
    args = parser.parse_args()
    project = args.project.expanduser().resolve()
    output = args.output.expanduser().resolve()
    if output.exists():
        print(f"Error: refusing to overwrite existing output: {output}", file=sys.stderr)
        return 2
    try:
        svgs = ordered_svgs(project)
        if not svgs:
            raise ValueError(f"No SVG files found in {project / 'svg_output'}")
        executor = str(detect_executor(args.executor)["name"])
        ppt_master = find_ppt_master()
        selected = args.backend
        if selected == "auto":
            if executor == "codex":
                raise RuntimeError(
                    "Codex detected: use the host Presentations skill/runtime for final "
                    "authoring and export. Pass --backend raster only to force the explicit "
                    "non-editable fallback."
                )
            selected = "ppt-master" if ppt_master else "raster"
        if selected == "ppt-master":
            if not ppt_master:
                message = "PPT Master not detected; set PPT_MASTER_ROOT or choose raster"
                if executor == "claude-code":
                    message += (
                        ". Claude Code needs PPT Master or another native PPTX backend before "
                        "editable output can be promised"
                    )
                raise RuntimeError(message)
            result = export_ppt_master(project, output, ppt_master)
            result["slide_count"] = len(svgs)
        else:
            if args.require_editable:
                raise RuntimeError("Editable export required, but only raster fallback is available")
            result = export_raster(project, output, svgs)
        result.update(
            {
                "schema": "teaching-deck-export.v1",
                "project": str(project),
                "output": str(output),
                "created_at": datetime.now(timezone.utc).isoformat(),
                "source_files": [path.name for path in svgs],
            }
        )
        report_path = args.report.expanduser().resolve() if args.report else output.with_suffix(".export.json")
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        print(json.dumps({"output": str(output), "backend": result["backend"], "editable": result["editable"], "slides": len(svgs)}))
        return 0
    except (OSError, ValueError, RuntimeError, subprocess.CalledProcessError, ET.ParseError) as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
